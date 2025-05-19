import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io

# Carregar o dataset
covid = pd.read_csv('dataset_covid.csv')

# Filtrar apenas janeiro de 2023
covid['Data'] = pd.to_datetime(covid['Data'])
covid_jan = covid[(covid['Data'] >= '2023-01-01') & (covid['Data'] <= '2023-01-31')]

# RQ1: Evolução dos casos acumulados por estado
cases_by_date_state = covid_jan.groupby(['Data', 'Estado'])['Casos_Acumulados'].sum().reset_index()

# Gráfico RQ1
plt.figure(figsize=(10,6))
for estado in cases_by_date_state['Estado'].unique():
    data_estado = cases_by_date_state[cases_by_date_state['Estado'] == estado]
    plt.plot(data_estado['Data'], data_estado['Casos_Acumulados'], label=estado)
plt.title('Evolução dos Casos Acumulados de COVID-19 por Estado (Jan/2023)')
plt.xlabel('Data')
plt.ylabel('Casos Acumulados')
plt.legend(bbox_to_anchor=(1.05, 1), loc='upper left', fontsize=8)
plt.tight_layout()
img_rq1 = io.BytesIO()
plt.savefig(img_rq1, format='png')
plt.close()
img_rq1.seek(0)

# RQ2: Taxa de letalidade média por estado em janeiro de 2023
letalidade_media = covid_jan.groupby('Estado')['Taxa_Letalidade'].mean().sort_values(ascending=False)

# Gráfico RQ2
plt.figure(figsize=(10,6))
letalidade_media.plot(kind='bar', color='tomato')
plt.title('Taxa de Letalidade Média por Estado (Jan/2023)')
plt.xlabel('Estado')
plt.ylabel('Taxa de Letalidade Média')
plt.tight_layout()
img_rq2 = io.BytesIO()
plt.savefig(img_rq2, format='png')
plt.close()
img_rq2.seek(0)

# Criar apresentação
prs = Presentation()

# Slide RQ1
title_slide_layout = prs.slide_layouts[5]
slide1 = prs.slides.add_slide(title_slide_layout)
title1 = slide1.shapes.title
if not title1:
    title1 = slide1.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(1)).text_frame
run = title1.text_frame.add_paragraph().add_run()
run.text = 'RQ1: Como evoluíram os casos acumulados de COVID-19 nos estados brasileiros em janeiro de 2023?'
run.font.size = Pt(20)
run.font.bold = True
slide1.shapes.add_picture(img_rq1, Inches(0.5), Inches(1.2), Inches(8.5), Inches(4.5))

# Explicação RQ1
left = Inches(0.5)
top = Inches(5.8)
width = Inches(8.5)
height = Inches(1.2)
text_box = slide1.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
tf.text = 'O gráfico mostra a evolução dos casos acumulados de COVID-19 em cada estado brasileiro durante janeiro de 2023. Cada linha representa um estado.'

# Slide RQ2
slide2 = prs.slides.add_slide(title_slide_layout)
title2 = slide2.shapes.title
if not title2:
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(9), Inches(1)).text_frame
run2 = title2.text_frame.add_paragraph().add_run()
run2.text = 'RQ2: Qual a taxa de letalidade média por estado em janeiro de 2023?'
run2.font.size = Pt(20)
run2.font.bold = True
slide2.shapes.add_picture(img_rq2, Inches(0.5), Inches(1.2), Inches(8.5), Inches(4.5))

# Explicação RQ2
text_box2 = slide2.shapes.add_textbox(left, top, width, height)
tf2 = text_box2.text_frame
tf2.text = 'O gráfico apresenta a taxa de letalidade média por estado no período analisado. A taxa de letalidade é a razão entre o número de óbitos e o número de casos.'

# Salvar apresentação
prs.save('Sprint2_RQs_COVID.pptx')
'Sprint2_RQs_COVID.pptx gerado com sucesso.'